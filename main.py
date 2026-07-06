"""
黒背景 PowerPoint → 白背景 PowerPoint 変換ツール (v1)

変換対象:
  - スライド背景を白に変換する
  - 白いテキストを黒に変換する（schemeClr:tx1 / prstClr:white / srgbClr近白）
  - 白い線・枠線・矢印を黒に変換する（同上）

使い方:
  python3 main.py sample1.pptx
  python3 main.py              # 引数なしは sample1.pptx を使用
"""

import colorsys
import os
import sys
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor

# OOXML のDrawingML名前空間
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# 白判定の閾値（各チャンネルが 240 以上なら白とみなす）
WHITE_THRESHOLD = 240
BLACK_HEX = '000000'


def is_white_hex(hex_val: str) -> bool:
    """6桁HEX文字列（例: 'FFFFFF'）が白っぽいかを判定する"""
    if len(hex_val) != 6:
        return False
    r = int(hex_val[0:2], 16)
    g = int(hex_val[2:4], 16)
    b = int(hex_val[4:6], 16)
    return r >= WHITE_THRESHOLD and g >= WHITE_THRESHOLD and b >= WHITE_THRESHOLD


def replace_with_black(parent: etree._Element, old: etree._Element) -> None:
    """solidFill 内の白色要素を黒の srgbClr に置き換える"""
    idx = list(parent).index(old)
    parent.remove(old)
    new = etree.Element(f'{{{NS_A}}}srgbClr')
    new.set('val', BLACK_HEX)
    parent.insert(idx, new)


def convert_solid_fill(sf: etree._Element) -> bool:
    """
    <a:solidFill> の色が白系なら黒に書き換える。変換したら True を返す。

    変換対象の色指定:
      srgbClr   : RGBが白に近い（例: FFFFFF, F2F2F2）
      schemeClr : val="tx1"（ダークテーマではlt1=白に解決される）
      prstClr   : val="white"（プリセット白）
      sysClr    : lastClr が白（システムカラー白）
    """
    for child in list(sf):
        tag = child.tag.replace(f'{{{NS_A}}}', '')

        if tag == 'srgbClr':
            if is_white_hex(child.get('val', '')):
                child.set('val', BLACK_HEX)
                return True

        elif tag == 'schemeClr':
            # tx1はダークテーマではlt1（白）に解決されるため黒に変換する
            if child.get('val') == 'tx1':
                replace_with_black(sf, child)
                return True

        elif tag == 'prstClr':
            if child.get('val') == 'white':
                replace_with_black(sf, child)
                return True

        elif tag == 'sysClr':
            # lastClr はシステム設定によらず実際に使われた色を示す
            if is_white_hex(child.get('lastClr', '')):
                replace_with_black(sf, child)
                return True

    return False


def apply_lum_adjust(base_rgb: tuple, lm: int, lo: int) -> tuple:
    """
    DrawingML の lumMod / lumOff を HLS 経由で RGB に適用する。
    base_rgb: (R, G, B) 各0-255
    lm: lumMod 値（例: 95000 = 95%）
    lo: lumOff 値（例: 0）
    """
    r, g, b = base_rgb[0] / 255, base_rgb[1] / 255, base_rgb[2] / 255
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    new_l = min(1.0, max(0.0, l * lm / 100000 + lo / 100000))
    r2, g2, b2 = colorsys.hls_to_rgb(h, new_l, s)
    return (round(r2 * 255), round(g2 * 255), round(b2 * 255))


def convert_shape_fill_sf(sf: etree._Element) -> bool:
    """
    図形塗りの <a:solidFill> を変換する。変換したら True を返す。

    【変換ルール】
    sysClr:window（システム白）: lumMod/lumOff を計算した明示的な srgbClr に変換する。
      例: window + lumMod=95000 → #F2F2F2（薄グレー）
          window + lumMod=50000 → #808080（中グレー）
    sysClr:windowText（システム黒）: srgbClr:000000 に変換する。

    【変換しないもの】
    srgbClr の直接指定: すでに明示的な色なので変換不要
    schemeClr: v2（clrMap 変更）で対応する
    """
    for child in list(sf):
        tag = child.tag.replace(f'{{{NS_A}}}', '')
        if tag != 'sysClr':
            continue

        val = child.get('val', '')
        last_clr = child.get('lastClr', 'FFFFFF')
        lm_e = child.find(f'{{{NS_A}}}lumMod')
        lo_e = child.find(f'{{{NS_A}}}lumOff')
        lm = int(lm_e.get('val', '100000')) if lm_e is not None else 100000
        lo = int(lo_e.get('val', '0'))       if lo_e is not None else 0

        if val == 'window':
            # sysClr:window = システム白(#FFFFFF)。lumMod/lumOff で実際のグレーを計算する
            rgb = apply_lum_adjust((255, 255, 255), lm, lo)
            hex_val = f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
        elif val == 'windowText':
            # sysClr:windowText = システム黒 → 明示的な黒に変換する
            hex_val = BLACK_HEX
        else:
            # その他の sysClr は lastClr（実績値）を使って変換する
            base = (int(last_clr[0:2], 16), int(last_clr[2:4], 16), int(last_clr[4:6], 16))
            rgb = apply_lum_adjust(base, lm, lo)
            hex_val = f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'

        idx = list(sf).index(child)
        sf.remove(child)
        new_e = etree.Element(f'{{{NS_A}}}srgbClr')
        new_e.set('val', hex_val)
        sf.insert(idx, new_e)
        return True

    return False


def convert_shape_fills(elem: etree._Element) -> int:
    """
    図形の塗り色（<p:spPr> 直下の <a:solidFill>）を変換する。
    変換した件数を返す。

    sysClr（システムカラー依存）のみを対象とし、明示的な srgbClr に変換する。
    srgbClr の直接指定や schemeClr はそのまま維持する（白→黒の一括変換は行わない）。
    """
    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    count = 0
    for sp_pr in elem.findall(f'{{{NS_P}}}spPr'):
        for sf in sp_pr.findall(f'{{{NS_A}}}solidFill'):
            if convert_shape_fill_sf(sf):
                count += 1
    return count


def convert_line_colors(elem: etree._Element) -> int:
    """
    要素内の <a:ln>（線・枠線・矢印）の白色を黒に変換する。
    変換した件数を返す。
    """
    count = 0
    for ln in elem.findall(f'.//{{{NS_A}}}ln'):
        for sf in ln.findall(f'{{{NS_A}}}solidFill'):
            if convert_solid_fill(sf):
                count += 1
    return count


def add_black_fill(rpr: etree._Element) -> None:
    """<a:rPr> に黒の solidFill を先頭に挿入する"""
    sf = etree.Element(f'{{{NS_A}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
    clr.set('val', BLACK_HEX)
    rpr.insert(0, sf)


def convert_text_colors(elem: etree._Element) -> int:
    """
    要素内のテキスト白色を黒に変換する。変換した件数を返す。

    対象1: <a:rPr> / <a:endParaRPr> に明示的な白色指定がある場合
    対象2: <a:r> にテキストがあるが solidFill がない場合
            → ダークテーマの tx1=white を継承しているため明示的に黒を設定する
    """
    count = 0

    # 対象1: 明示的な白色をもつ rPr と endParaRPr を変換する
    for tag in ('rPr', 'endParaRPr'):
        for rpr in elem.findall(f'.//{{{NS_A}}}{tag}'):
            for sf in rpr.findall(f'{{{NS_A}}}solidFill'):
                if convert_solid_fill(sf):
                    count += 1

    # 対象2: 色指定のないランにテーマ継承の白が適用されているため明示的に黒を設定する
    for run in elem.findall(f'.//{{{NS_A}}}r'):
        t_elem = run.find(f'{{{NS_A}}}t')
        if t_elem is None or not (t_elem.text or '').strip():
            continue  # テキストのないランはスキップ

        rpr = run.find(f'{{{NS_A}}}rPr')
        if rpr is None:
            # rPr 自体がない場合は新規作成して黒を設定する
            rpr = etree.Element(f'{{{NS_A}}}rPr')
            run.insert(0, rpr)
            add_black_fill(rpr)
            count += 1
        elif rpr.find(f'{{{NS_A}}}solidFill') is None:
            # rPr はあるが色指定がない場合も黒を追加する
            add_black_fill(rpr)
            count += 1

    return count


def convert_slide(slide, slide_num: int) -> dict:
    """1枚のスライドに背景・テキスト・線の変換を適用する"""
    result = {'slide': slide_num, 'text': 0, 'line': 0, 'fill': 0}

    # 背景を白い単色塗りに変換する
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for shape in slide.shapes:
        elem = shape._element
        result['text'] += convert_text_colors(elem)
        result['line'] += convert_line_colors(elem)
        result['fill'] += convert_shape_fills(elem)

    return result


def convert_pptx(input_path: str) -> str:
    """pptxを読み込み変換して保存する。保存先のパスを返す。"""
    prs = Presentation(input_path)

    print(f"変換開始: {input_path}  (スライド数: {len(prs.slides)})")

    total_text = 0
    total_line = 0

    total_fill = 0
    for i, slide in enumerate(prs.slides, start=1):
        result = convert_slide(slide, i)
        total_text += result['text']
        total_line += result['line']
        total_fill += result['fill']
        print(f"  スライド{i}: 文字 {result['text']} 件、線 {result['line']} 件、図形塗り {result['fill']} 件を変換")

    # 元ファイルを上書きせず converted_ プレフィックスで保存する
    dir_name = os.path.dirname(os.path.abspath(input_path))
    base_name = os.path.basename(input_path)
    output_path = os.path.join(dir_name, f"converted_{base_name}")

    prs.save(output_path)

    print(f"\n合計: 文字 {total_text} 件、線 {total_line} 件、図形塗り {total_fill} 件を変換")
    print(f"保存先: {output_path}")
    return output_path


if __name__ == "__main__":
    input_file = sys.argv[1] if len(sys.argv) > 1 else "sample1.pptx"

    if not os.path.exists(input_file):
        print(f"エラー: '{input_file}' が見つかりません")
        sys.exit(1)

    convert_pptx(input_file)
